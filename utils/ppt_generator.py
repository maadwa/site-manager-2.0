from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import plotly.graph_objects as go
from typing import List, Dict, Any
import os
import tempfile
import datetime

class PPTGenerator:
    """Handles PowerPoint presentation generation"""
    
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
    
    def create_presentation(self, graphs: List[Dict[str, Any]], title: str = "Construction Project Report") -> str:
        """Create PowerPoint presentation with graphs"""
        try:
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            self._add_title_slide(prs, title)
            
            # Add slides for each graph
            for i, graph_data in enumerate(graphs):
                self._add_graph_slide(prs, graph_data, i + 1)
            
            # Add summary slide
            self._add_summary_slide(prs, len(graphs))
            
            # Save presentation
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
            filepath = os.path.join(self.temp_dir, filename)
            
            prs.save(filepath)
            
            return filepath
        
        except Exception as e:
            print(f"Error creating presentation: {str(e)}")
            return None
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """Add title slide to presentation"""
        try:
            # Use title slide layout
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Set title
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            
            # Set subtitle
            subtitle_placeholder = slide.placeholders[1]
            current_date = datetime.datetime.now().strftime("%B %d, %Y")
            subtitle_placeholder.text = f"Generated on {current_date}\nConstruction Project Management Dashboard"
            
            # Format title
            title_paragraph = title_placeholder.text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(44)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(31, 73, 125)
            
            # Format subtitle
            for paragraph in subtitle_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(68, 84, 106)
        
        except Exception as e:
            print(f"Error adding title slide: {str(e)}")
    
    def _add_graph_slide(self, prs: Presentation, graph_data: Dict[str, Any], slide_number: int):
        """Add slide with graph"""
        try:
            # Use content slide layout
            content_slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(1)
            )
            title_frame = title_shape.text_frame
            title_frame.text = graph_data.get('title', f'Chart {slide_number}')
            
            # Format title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(28)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(31, 73, 125)
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            # Add graph image
            if 'figure' in graph_data:
                fig = graph_data['figure']
                
                # Save figure as temporary image
                img_path = os.path.join(self.temp_dir, f"chart_{slide_number}.png")
                fig.write_image(img_path, width=800, height=500, scale=2)
                
                # Add image to slide
                slide.shapes.add_picture(
                    img_path, 
                    Inches(1), Inches(2), 
                    Inches(8), Inches(5)
                )
            
            # Add slide number
            slide_number_shape = slide.shapes.add_textbox(
                Inches(8.5), Inches(7), Inches(1), Inches(0.5)
            )
            slide_number_frame = slide_number_shape.text_frame
            slide_number_frame.text = str(slide_number)
            slide_number_paragraph = slide_number_frame.paragraphs[0]
            slide_number_paragraph.font.size = Pt(12)
            slide_number_paragraph.font.color.rgb = RGBColor(127, 127, 127)
            slide_number_paragraph.alignment = PP_ALIGN.RIGHT
        
        except Exception as e:
            print(f"Error adding graph slide: {str(e)}")
    
    def _add_summary_slide(self, prs: Presentation, graph_count: int):
        """Add summary slide"""
        try:
            # Use content slide layout
            content_slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set title
            title_placeholder = slide.shapes.title
            title_placeholder.text = "Report Summary"
            
            # Format title
            title_paragraph = title_placeholder.text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(36)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(31, 73, 125)
            
            # Add content
            content_placeholder = slide.placeholders[1]
            content_frame = content_placeholder.text_frame
            content_frame.clear()
            
            # Add summary points
            summary_points = [
                f"Total visualizations generated: {graph_count}",
                "Data sourced from construction project Excel files",
                "Graphs created based on admin-selected columns",
                "Real-time data analysis and visualization",
                "Export capability for project reporting"
            ]
            
            for point in summary_points:
                p = content_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(68, 84, 106)
                p.space_after = Pt(12)
            
            # Add footer
            footer_shape = slide.shapes.add_textbox(
                Inches(1), Inches(6.5), Inches(8), Inches(1)
            )
            footer_frame = footer_shape.text_frame
            footer_frame.text = "Generated by Construction Project Management Dashboard"
            footer_paragraph = footer_frame.paragraphs[0]
            footer_paragraph.font.size = Pt(12)
            footer_paragraph.font.italic = True
            footer_paragraph.font.color.rgb = RGBColor(127, 127, 127)
            footer_paragraph.alignment = PP_ALIGN.CENTER
        
        except Exception as e:
            print(f"Error adding summary slide: {str(e)}")
    
    def add_data_slide(self, prs: Presentation, data_summary: Dict[str, Any]):
        """Add slide with data summary"""
        try:
            # Use content slide layout
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set title
            title_placeholder = slide.shapes.title
            title_placeholder.text = "Data Summary"
            
            # Add content
            content_placeholder = slide.placeholders[1]
            content_frame = content_placeholder.text_frame
            content_frame.clear()
            
            # Add data summary
            for key, value in data_summary.items():
                p = content_frame.add_paragraph()
                p.text = f"{key}: {value}"
                p.font.size = Pt(16)
                p.space_after = Pt(8)
        
        except Exception as e:
            print(f"Error adding data slide: {str(e)}")
    
    def cleanup_temp_files(self):
        """Clean up temporary files"""
        try:
            import shutil
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
        except Exception as e:
            print(f"Error cleaning up temp files: {str(e)}")
